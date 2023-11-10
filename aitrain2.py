import streamlit as st
from pptx import Presentation
import os
from github import Github
import openai
import difflib
import openpyxl
from docx import Document
import fitz

st.set_page_config(page_title="Chat with the Bain Report", page_icon="🦙", layout="centered", initial_sidebar_state="auto", menu_items=None)

if "OPENAI_API_KEY" not in st.secrets or "GITHUB_TOKEN" not in st.secrets:
    st.error("Please set the necessary secrets (OPENAI_API_KEY and GITHUB_TOKEN) on the Streamlit dashboard.")
    sys.exit(1)

openai.api_key = st.secrets["OPENAI_API_KEY"]
github_token = st.secrets["GITHUB_TOKEN"]

g = Github(github_token)
repo = g.get_repo("scooter7/aitrain")

def upload_to_github(file_path, repo, path_in_repo):
    with open(file_path, "rb") as file:
        content = file.read()
    git_file = path_in_repo + '/' + os.path.basename(file_path)
    try:
        contents = repo.get_contents(git_file)
        repo.update_file(contents.path, "Updating file", content, contents.sha)
        st.success('File updated on GitHub')
    except:
        repo.create_file(git_file, "Creating new file", content)
        st.success('File created on GitHub')

def save_uploaded_file(uploaded_file):
    with open(os.path.join("tempDir", uploaded_file.name), "wb") as f:
        f.write(uploaded_file.getbuffer())
    return os.path.join("tempDir", uploaded_file.name)

def get_document_titles_and_urls(repo):
    contents = repo.get_contents("docs")
    document_titles = []
    document_urls = {}
    for content_file in contents:
        if content_file.type == "file" and content_file.name.endswith(('.pdf', '.docx', '.xlsx', '.pptx')):
            document_titles.append(content_file.name)
            document_urls[content_file.name] = content_file.download_url
    return document_titles, document_urls

def extract_text_by_stages(pptx_path):
    prs = Presentation(pptx_path)
    stages_text = {
        "Stage 1": [2, 9],
        "Stage 2": [10, 14],
        "Stage 3": [15, 21],
        "Stage 4": [22, 24],
        "Stage 5": [25, 30],
        "Stage 6": [31, 36],
    }
    stages_content = {}
    for stage, (start_slide, end_slide) in stages_text.items():
        text_content = []
        for slide number in range(start_slide - 1, end_slide):
            slide = prs.slides[slide_number]
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text.strip())
        stages_content[stage] = "\n".join(text_content)
    return stages_content

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def extract_text_from_pdf(pdf_path):
    with fitz.open(pdf_path) as doc:
        text = ""
        for page in doc:
            text += page.get_text()
    return text

def extract_data_from_xlsx(xlsx_path):
    workbook = openpyxl.load_workbook(xlsx_path)
    text = ""
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        for row in worksheet.iter_rows(values_only=True):
            text += " ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    return text

document_titles, document_urls = get_document_titles_and_urls(repo)

stages_content = extract_text_by_stages("docs/marketing_strategy_plan_methodology.pptx")

if 'current_stage_index' not in st.session_state:
    st.session_state.current_stage_index = 0

current_stage_keys = list(stages_content.keys())

if st.button("Go to next stage"):
    st.session_state.current_stage_index += 1
    if st.session_state.current_stage_index >= len(current_stage_keys):
        st.session_state.current_stage_index = 0

current_stage = current_stage_keys[st.session_state.current_stage_index]
st.subheader(current_stage)
st.write(stages_content[current_stage])

uploaded_file = st.file_uploader("Upload your document", type=['docx', 'xlsx', 'pptx', 'pdf'])
if uploaded_file is not None:
    file_path = save_uploaded_file(uploaded_file)
    upload_to_github(file_path, repo, "uploads")
    file_content = ""
    if uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        file_content = extract_text_from_docx(file_path)
    elif uploaded_file.type == "application/pdf":
        file_content = extract_text_from_pdf(file_path)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        file_content = extract_data_from_xlsx(file_path)
    st.session_state.uploaded_file_content = file_content

if "messages" not in st.session_state:
    st.session_state.messages = []

if prompt := st.text_input("Your question"):
    st.session_state.messages.append({"role": "user", "content": prompt})

for message in st.session_state.messages:
    st.write(f"{message['role'].title()}: {message['content']}")

if st.session_state.messages and st.session_state.messages[-1]["role"] != "assistant":
    document_keywords = ['document', 'file', 'download', 'link', 'template', 'worksheet', 'form']
    if any(keyword in prompt.lower() for keyword in document_keywords):
        closest_matches = difflib.get_close_matches(prompt.lower(), [title.lower() for title in document_titles], n=5, cutoff=0.3)
        if closest_matches:
            response_content = "Here are the documents that might match your request:\n"
            for title in closest_matches:
                document_url = document_urls[title]
                response_content += f"- [{title}]({document_url})\n"
        else:
            response_content = "I couldn't find the document you're looking for. Please make sure to use the exact title of the document or provide more context."
    else:
        context = st.session_state.uploaded_file_content[:1000] if 'uploaded_file_content' in st.session_state else ""
        response = openai.chat.completions.create(
            model="gpt-4",
            messages=[{"role": "system", "content": "You are a helpful assistant."}] + st.session_state.messages
        )
        response_content = response.choices[0].message.content
        st.session_state.messages.append({"role": "assistant", "content": response_content})
        st.write(f"Assistant: {response_content}")
