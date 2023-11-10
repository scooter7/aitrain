import streamlit as st
from pptx import Presentation
import os
from github import Github
import openai
import difflib
import openpyxl
from docx import Document
import fitz

# Set Streamlit page configuration
st.set_page_config(page_title="Chat with the Bain Report", page_icon="🦙", layout="centered", initial_sidebar_state="auto")

# Check for necessary secrets
if "OPENAI_API_KEY" not in st.secrets or "GITHUB_TOKEN" not in st.secrets:
    st.error("Please set the necessary secrets (OPENAI_API_KEY and GITHUB_TOKEN) on the Streamlit dashboard.")
    sys.exit(1)

# Set API keys
openai.api_key = st.secrets["OPENAI_API_KEY"]
github_token = st.secrets["GITHUB_TOKEN"]

# Initialize GitHub
g = Github(github_token)
repo = g.get_repo("scooter7/aitrain")

# Function to upload file to GitHub
def upload_to_github(file_path, repo, path_in_repo):
    with open(file_path, "rb") as file:
        content = file.read()
    git_file = path_in_repo + '/' + os.path.basename(file_path)
    try:
        contents = repo.get_contents(git_file)
        repo.update_file(contents.path, "Updating file", content, contents.sha)
        st.success('File updated on GitHub')
    except:
        repo.create_file(git_file, "Creating new file", content)
        st.success('File created on GitHub')

# Function to save uploaded file
def save_uploaded_file(uploaded_file):
    file_path = os.path.join("uploads", uploaded_file.name)
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    return file_path

# Function to get document titles and URLs from the docs directory in the repository
def get_document_titles_and_urls(repo):
    contents = repo.get_contents("docs")  # Gets the contents of the 'docs' directory
    document_titles = []
    document_urls = {}

    for content_file in contents:
        if content_file.type == "file" and content_file.name.endswith(('.pdf', '.docx', '.xlsx')):
            document_titles.append(content_file.name)
            # Generate GitHub URL for the file
            document_url = f"https://github.com/{repo.full_name}/blob/main/docs/{content_file.name}"
            document_urls[content_file.name] = document_url

    return document_titles, document_urls

# Fetch the actual document titles and URLs
document_titles, document_urls = get_document_titles_and_urls(repo)

# Function to extract text by stages from a PowerPoint presentation
def extract_text_by_stages(pptx_path):
    prs = Presentation(pptx_path)
    stages_text = {
        "Stage 1": [2, 9],
        "Stage 2": [10, 14],
        "Stage 3": [15, 21],
        "Stage 4": [22, 24],
        "Stage 5": [25, 30],
        "Stage 6": [31, 36],
    }
    stages_content = {}
    for stage, (start_slide, end_slide) in stages_text.items():
        text_content = []
        for slide_number in range(start_slide - 1, end_slide):
            slide = prs.slides[slide_number]
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text.strip())
        stages_content[stage] = "\n".join(text_content)
    return stages_content

# Function to find action items in a stage
def find_action_items_in_stage(stage_content):
    action_items = {}
    lines = stage_content.split("\n")
    for line in lines:
        if "Action Item –" in line:
            action_item_title = line.split("–")[1].strip()
            action_items[action_item_title] = None
    return action_items

# Function to map action items to files
def map_action_items_to_files(action_items, document_titles, document_urls):
    for item in action_items.keys():
        closest_match = difflib.get_close_matches(item, document_titles, n=1, cutoff=0.5)
        if closest_match:
            action_items[item] = document_urls[closest_match[0]]

# Function to extract text from a DOCX file
def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

# Function to extract data from an XLSX file
def extract_data_from_xlsx(xlsx_path):
    workbook = openpyxl.load_workbook(xlsx_path)
    text = ""
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        for row in worksheet.iter_rows(values_only=True):
            text += " ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    return text

# Function to summarize text
def summarize_text(text, max_length=500):
    if len(text) > max_length:
        return text[:max_length] + "..."
    return text

# Main code
document_titles, document_urls = get_document_titles_and_github_urls(repo)
stages_content = extract_text_by_stages("docs/marketing_strategy_plan_methodology.pptx")

if 'current_stage_index' not in st.session_state:
    st.session_state.current_stage_index = 0

current_stage_keys = list(stages_content.keys())

if st.button("Go to next stage"):
    st.session_state.current_stage_index += 1
    if st.session_state.current_stage_index >= len(current_stage_keys):
        st.session_state.current_stage_index = 0

current_stage = current_stage_keys[st.session_state.current_stage_index]
current_stage_content = stages_content[current_stage]
st.subheader(current_stage)
st.write(current_stage_content)

action_items = find_action_items_in_stage(current_stage_content)
map_action_items_to_files(action_items, document_titles, document_urls)
for title, url in action_items.items():
    st.markdown(f"- [{title}]({url})")

uploaded_file = st.file_uploader("Upload your document", type=['docx', 'xlsx', 'pptx', 'pdf'])
if uploaded_file is not None:
    file_path = save_uploaded_file(uploaded_file)
    upload_to_github(file_path, repo, "uploads")
    file_content = ""
    if uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        file_content = extract_text_from_docx(file_path)
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        file_content = extract_data_from_xlsx(file_path)
    summarized_content = summarize_text(file_content)
    st.write("Summary of uploaded document:")
    st.write(summarized_content)

if "messages" not in st.session_state:
    st.session_state.messages = []

if prompt := st.text_input("Your question"):
    st.session_state.messages.append({"role": "user", "content": prompt})

for message in st.session_state.messages:
    st.write(f"{message['role'].title()}: {message['content']}")

if st.session_state.messages and st.session_state.messages[-1]["role"] != "assistant":
    response = openai.chat.completions.create(
        model="gpt-4",
        messages=st.session_state.messages
    )
    response_content = response.choices[0].message.content
    st.session_state.messages.append({"role": "assistant", "content": response_content})
    st.write(f"Assistant: {response_content}")
